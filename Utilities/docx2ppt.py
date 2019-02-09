import docx
from pptx import Presentation

doc=docx.Document('FLFinalSPv1.docx')
li=[]
a=0
while a < len(doc.paragraphs):
    li2=[]
    if (doc.paragraphs[a].text).isupper():
        li2.append(doc.paragraphs[a].text)
        a=a+1
        while ((doc.paragraphs[a].text).isupper())== False:
            li2.append(doc.paragraphs[a].text)
            a=a+1
    if len(li2) > 0:
        li.append(li2)
    a=a+1


m1=[]
for i in li:
    m2=[]
    for j in i:
        if len(j)>0:
            m2.append(j)
    m1.append(m2)



prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
for i in range(len(m1)):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    #subtitle = slide.placeholders[1]

    title.text = '\n'.join(m1[i])
    

prs.save('output_ppt.pptx')
